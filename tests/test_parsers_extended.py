"""
اختبارات إضافية موسعة لمفسرات ملفات DOCX, PPTX, XLSX, PDF
"""

import io
import pytest
import docx
import pptx
import openpyxl
import pypdf

from document_parser import DocumentParser


def test_docx_parser():
    """اختبار استخراج النصوص من مستندات Word"""
    doc = docx.Document()
    doc.add_heading("مقدمة في الذكاء الاصطناعي", level=1)
    doc.add_paragraph("الذكاء الاصطناعي هو فرع من علوم الحاسوب يهدف لبناء أنظمة ذكية.")
    
    # جدول
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "المفهوم"
    table.cell(0, 1).text = "الوصف"
    table.cell(1, 0).text = "RAG"
    table.cell(1, 1).text = "توليد معزز بالاسترجاع"

    stream = io.BytesIO()
    doc.save(stream)
    docx_bytes = stream.getvalue()

    sections = DocumentParser.extract_text("ai_intro.docx", docx_bytes)
    assert len(sections) >= 1
    combined = " ".join(s.text for s in sections)
    assert "مقدمة في الذكاء الاصطناعي" in combined
    assert "توليد معزز بالاسترجاع" in combined


def test_pptx_parser():
    """اختبار استخراج النصوص من عروض PowerPoint"""
    prs = pptx.Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    txBox = slide.shapes.add_textbox(0, 0, 100, 100)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "شريحة المحاضرة الأولى: مفاهيم نظم المعلومات"

    stream = io.BytesIO()
    prs.save(stream)
    pptx_bytes = stream.getvalue()

    sections = DocumentParser.extract_text("lecture.pptx", pptx_bytes)
    assert len(sections) == 1
    assert "نظم المعلومات" in sections[0].text
    assert sections[0].section_type == "شريحة"
    assert sections[0].section_number == 1


def test_xlsx_parser():
    """اختبار استخراج النصوص من جداول Excel"""
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "درجات الطلاب"
    ws.append(["الاسم", "المشروع", "الدرجة النهائية"])
    ws.append(["محمد", "مُجاوِب", "100"])

    stream = io.BytesIO()
    wb.save(stream)
    xlsx_bytes = stream.getvalue()

    sections = DocumentParser.extract_text("students.xlsx", xlsx_bytes)
    assert len(sections) == 1
    assert "درجات الطلاب" in sections[0].text
    assert "مُجاوِب" in sections[0].text
    assert "100" in sections[0].text


def test_pdf_parser():
    """اختبار استخراج النصوص من ملف PDF"""
    writer = pypdf.PdfWriter()
    page = writer.add_blank_page(width=200, height=200)
    
    stream = io.BytesIO()
    writer.write(stream)
    pdf_bytes = stream.getvalue()

    sections = DocumentParser.extract_text("empty_test.pdf", pdf_bytes)
    assert len(sections) >= 1
