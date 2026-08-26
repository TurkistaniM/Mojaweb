"""
وحدة معالجة واستخراج النصوص من الملفات المتعددة (Document Parser)
تدعم:
- المستندات: PDF, DOCX, PPTX, XLSX, CSV, TXT, MD
- الصوتيات: MP3, WAV, M4A, OGG, AAC, FLAC
- الفيديو: MP4, MOV, WEBM, M4V
"""

import io
import csv
import os
import tempfile
from typing import List, Dict, Any, Tuple
from dataclasses import dataclass

# استيراد مفسرات المستندات
import pypdf
import docx
import pptx
import openpyxl

# مفسرات الصوت والفيديو
import speech_recognition as sr
from pydub import AudioSegment


@dataclass
class ExtractedSection:
    filename: str
    section_number: int  # رقم الصفحة / الشريحة / المقطع الزمني
    section_type: str    # "صفحة", "شريحة", "ورقة عمل", "مقطع", "مقطع صوتي/مرئي"
    text: str

    @property
    def source_label(self) -> str:
        return f"{self.filename} ({self.section_type} {self.section_number})"


class DocumentParser:
    """مستخرج النصوص والمحتوى من المستندات والصوتيات ومقاطع الفيديو."""

    DOCUMENT_EXTENSIONS = {".pdf", ".docx", ".pptx", ".txt", ".md", ".csv", ".xlsx"}
    AUDIO_EXTENSIONS = {".mp3", ".wav", ".m4a", ".ogg", ".aac", ".flac", ".wma"}
    VIDEO_EXTENSIONS = {".mp4", ".mov", ".webm", ".m4v", ".mkv", ".avi"}

    SUPPORTED_EXTENSIONS = DOCUMENT_EXTENSIONS | AUDIO_EXTENSIONS | VIDEO_EXTENSIONS

    @classmethod
    def is_supported(cls, filename: str) -> bool:
        ext = os.path.splitext(filename.lower())[1]
        return ext in cls.SUPPORTED_EXTENSIONS

    @classmethod
    def is_audio_or_video(cls, filename: str) -> bool:
        ext = os.path.splitext(filename.lower())[1]
        return ext in (cls.AUDIO_EXTENSIONS | cls.VIDEO_EXTENSIONS)

    @classmethod
    def extract_text(cls, filename: str, file_bytes: bytes) -> List[ExtractedSection]:
        """
        استخراج النصوص من محتوى الملف مع الحفاظ على أرقام الصفحات/الشرائح أو التوقيت الزمني
        """
        ext = os.path.splitext(filename.lower())[1]

        if ext == ".pdf":
            return cls._extract_pdf(filename, file_bytes)
        elif ext == ".docx":
            return cls._extract_docx(filename, file_bytes)
        elif ext == ".pptx":
            return cls._extract_pptx(filename, file_bytes)
        elif ext == ".xlsx":
            return cls._extract_xlsx(filename, file_bytes)
        elif ext == ".csv":
            return cls._extract_csv(filename, file_bytes)
        elif ext in {".txt", ".md"}:
            return cls._extract_plain_text(filename, file_bytes)
        elif ext in cls.AUDIO_EXTENSIONS or ext in cls.VIDEO_EXTENSIONS:
            return cls._extract_audio_video(filename, file_bytes, ext)
        else:
            raise ValueError(f"امتداد الملف غير مدعوم: {ext}")

    @staticmethod
    def _extract_pdf(filename: str, file_bytes: bytes) -> List[ExtractedSection]:
        sections = []
        stream = io.BytesIO(file_bytes)
        reader = pypdf.PdfReader(stream)

        for i, page in enumerate(reader.pages, start=1):
            text = page.extract_text() or ""
            text = text.strip()
            if text:
                sections.append(ExtractedSection(
                    filename=filename,
                    section_number=i,
                    section_type="صفحة",
                    text=text
                ))

        if not sections:
            sections.append(ExtractedSection(
                filename=filename,
                section_number=1,
                section_type="صفحة",
                text="[ملف PDF فارغ أو يحتوي على صور ممسوحة ضوئياً فقط]"
            ))
        return sections

    @staticmethod
    def _extract_docx(filename: str, file_bytes: bytes) -> List[ExtractedSection]:
        sections = []
        stream = io.BytesIO(file_bytes)
        doc = docx.Document(stream)

        full_text_blocks = []
        for p in doc.paragraphs:
            t = p.text.strip()
            if t:
                full_text_blocks.append(t)

        for table in doc.tables:
            table_rows = []
            for row in table.rows:
                row_cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_cells:
                    table_rows.append(" | ".join(row_cells))
            if table_rows:
                full_text_blocks.append("\n".join(table_rows))

        combined_text = "\n\n".join(full_text_blocks).strip()
        if not combined_text:
            combined_text = "[مستند Word فارغ]"

        chunk_size = 1500
        paragraphs = combined_text.split("\n\n")
        current_chunk = []
        current_len = 0
        section_idx = 1

        for para in paragraphs:
            current_chunk.append(para)
            current_len += len(para)
            if current_len >= chunk_size:
                sections.append(ExtractedSection(
                    filename=filename,
                    section_number=section_idx,
                    section_type="مقطع",
                    text="\n\n".join(current_chunk)
                ))
                section_idx += 1
                current_chunk = []
                current_len = 0

        if current_chunk:
            sections.append(ExtractedSection(
                filename=filename,
                section_number=section_idx,
                section_type="مقطع",
                text="\n\n".join(current_chunk)
            ))

        return sections

    @staticmethod
    def _extract_pptx(filename: str, file_bytes: bytes) -> List[ExtractedSection]:
        sections = []
        stream = io.BytesIO(file_bytes)
        presentation = pptx.Presentation(stream)

        for i, slide in enumerate(presentation.slides, start=1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        t = paragraph.text.strip()
                        if t:
                            slide_texts.append(t)
                elif shape.has_table:
                    for row in shape.table.rows:
                        row_cells = [c.text.strip() for c in row.cells if c.text.strip()]
                        if row_cells:
                            slide_texts.append(" | ".join(row_cells))

            text = "\n".join(slide_texts).strip()
            if text:
                sections.append(ExtractedSection(
                    filename=filename,
                    section_number=i,
                    section_type="شريحة",
                    text=text
                ))

        if not sections:
            sections.append(ExtractedSection(
                filename=filename,
                section_number=1,
                section_type="شريحة",
                text="[عرض تقديمي فارغ أو بدون نصوص]"
            ))
        return sections

    @staticmethod
    def _extract_xlsx(filename: str, file_bytes: bytes) -> List[ExtractedSection]:
        sections = []
        stream = io.BytesIO(file_bytes)
        wb = openpyxl.load_workbook(stream, data_only=True)

        for i, sheetname in enumerate(wb.sheetnames, start=1):
            sheet = wb[sheetname]
            sheet_rows = []
            for row in sheet.iter_rows(values_only=True):
                non_empty = [str(val).strip() for val in row if val is not None and str(val).strip()]
                if non_empty:
                    sheet_rows.append(" | ".join(non_empty))

            text = f"ورقة: {sheetname}\n" + "\n".join(sheet_rows)
            text = text.strip()
            if sheet_rows:
                sections.append(ExtractedSection(
                    filename=filename,
                    section_number=i,
                    section_type="ورقة عمل",
                    text=text
                ))

        if not sections:
            sections.append(ExtractedSection(
                filename=filename,
                section_number=1,
                section_type="ورقة عمل",
                text="[ملف Excel فارغ]"
            ))
        return sections

    @staticmethod
    def _extract_csv(filename: str, file_bytes: bytes) -> List[ExtractedSection]:
        text_content = ""
        for encoding in ["utf-8-sig", "utf-8", "windows-1256", "latin-1"]:
            try:
                text_content = file_bytes.decode(encoding)
                break
            except UnicodeDecodeError:
                continue

        reader = csv.reader(io.StringIO(text_content))
        rows = []
        for r in reader:
            clean_cells = [c.strip() for c in r if c.strip()]
            if clean_cells:
                rows.append(" | ".join(clean_cells))

        combined = "\n".join(rows).strip()
        if not combined:
            combined = "[ملف CSV فارغ]"

        return [ExtractedSection(
            filename=filename,
            section_number=1,
            section_type="جدول",
            text=combined
        )]

    @staticmethod
    def _extract_plain_text(filename: str, file_bytes: bytes) -> List[ExtractedSection]:
        text_content = ""
        for encoding in ["utf-8-sig", "utf-8", "windows-1256", "latin-1"]:
            try:
                text_content = file_bytes.decode(encoding)
                break
            except UnicodeDecodeError:
                continue

        text_content = text_content.strip()
        if not text_content:
            text_content = "[ملف نصي فارغ]"

        chunk_size = 1500
        paragraphs = text_content.split("\n\n")
        sections = []
        current_chunk = []
        current_len = 0
        section_idx = 1

        for para in paragraphs:
            current_chunk.append(para)
            current_len += len(para)
            if current_len >= chunk_size:
                sections.append(ExtractedSection(
                    filename=filename,
                    section_number=section_idx,
                    section_type="مقطع",
                    text="\n\n".join(current_chunk)
                ))
                section_idx += 1
                current_chunk = []
                current_len = 0

        if current_chunk:
            sections.append(ExtractedSection(
                filename=filename,
                section_number=section_idx,
                section_type="مقطع",
                text="\n\n".join(current_chunk)
            ))

        return sections

    @classmethod
    def _extract_audio_video(cls, filename: str, file_bytes: bytes, ext: str) -> List[ExtractedSection]:
        """
        تفريغ الصوت ومقاطع الفيديو إلى نصوص مقسمة زمنياً
        Speech-to-Text Transcription with Timestamps
        """
        sections = []
        recognizer = sr.Recognizer()
        recognizer.energy_threshold = 300
        recognizer.dynamic_energy_threshold = True

        # حفظ الملف مؤقتاً للمعالجة
        suffix = ext.lower()
        with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as temp_in:
            temp_in.write(file_bytes)
            temp_in_path = temp_in.name

        temp_wav_path = None
        try:
            # تحويل المقطع الصوتي/المرئي إلى صيغة WAV أحادية القناة 16kHz
            format_name = ext.lstrip(".").lower()
            if format_name in {"mp4", "mov", "m4v", "mkv", "webm", "avi"}:
                # استخراج الصوت من ملف الفيديو
                audio = AudioSegment.from_file(temp_in_path)
            else:
                audio = AudioSegment.from_file(temp_in_path, format=format_name)

            audio = audio.set_channels(1).set_frame_rate(16000)

            # تجزئة الصوت إلى مقاطع بطول 30 ثانية لكل مقطع
            segment_len_ms = 30 * 1000
            total_segments = len(audio) // segment_len_ms + (1 if len(audio) % segment_len_ms > 0 else 0)

            for seg_idx in range(total_segments):
                start_ms = seg_idx * segment_len_ms
                end_ms = min(len(audio), (seg_idx + 1) * segment_len_ms)
                seg_audio = audio[start_ms:end_ms]

                # تنسيق الوقت (دقائق:ثوانٍ)
                start_sec = start_ms // 1000
                end_sec = end_ms // 1000
                time_label = f"{start_sec//60:02d}:{start_sec%60:02d} - {end_sec//60:02d}:{end_sec%60:02d}"

                # تفريغ المقطع
                seg_stream = io.BytesIO()
                seg_audio.export(seg_stream, format="wav")
                seg_stream.seek(0)

                with sr.AudioFile(seg_stream) as source:
                    audio_data = recognizer.record(source)
                    transcript = ""
                    try:
                        # محاولة تفريغ الكلام باللغة العربية أولاً ثم الإنجليزية
                        try:
                            transcript = recognizer.recognize_google(audio_data, language="ar-SA")
                        except Exception:
                            transcript = recognizer.recognize_google(audio_data, language="en-US")
                    except Exception:
                        transcript = ""

                    if transcript.strip():
                        sections.append(ExtractedSection(
                            filename=filename,
                            section_number=seg_idx + 1,
                            section_type=f"توقيت {time_label}",
                            text=f"[توقيت: {time_label}]\n{transcript.strip()}"
                        ))

            if not sections:
                duration_sec = len(audio) // 1000
                sections.append(ExtractedSection(
                    filename=filename,
                    section_number=1,
                    section_type="تسجيل",
                    text=f"[ملف صوتي/مرئي: {filename} - المدة: {duration_sec//60:02d}:{duration_sec%60:02d} دقيقة - لم يتم رصد كلام واضح أو تعذر الاتصال بخدمة التفريغ]"
                ))

        except Exception as e:
            # في حال تعذر تحويل الصوت بسبب عدم توفر ffmpeg محلياً، نوفر مقطعاً توثيقياً لبيانات الملف
            file_mb = len(file_bytes) / (1024 * 1024)
            sections.append(ExtractedSection(
                filename=filename,
                section_number=1,
                section_type="تسجيل",
                text=f"[ملف تسجيل {filename} بحجم {file_mb:.1f}MB - تم استلام الملف بنجاح للتسجيل الدراسي: {str(e)}]"
            ))
        finally:
            if os.path.exists(temp_in_path):
                try:
                    os.remove(temp_in_path)
                except Exception:
                    pass

        return sections
